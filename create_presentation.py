"""
Enterprise Pitch Deck Generator for AI MCP Gateway Solution
Generates a 16:9 widescreen PowerPoint presentation (MCP_Gateway_Enterprise_Solution.pptx)
with modern design, dark enterprise styling, structured card layouts, and visual hierarchy.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_PPTX = os.path.join(os.path.dirname(os.path.abspath(__file__)), "MCP_Gateway_Enterprise_Solution.pptx")

# Enterprise Color Palette
COLOR_BG = RGBColor(15, 23, 42)          # Slate 900
COLOR_CARD_BG = RGBColor(30, 41, 59)     # Slate 800
COLOR_CARD_BORDER = RGBColor(51, 65, 85) # Slate 700
COLOR_ACCENT = RGBColor(56, 189, 248)    # Sky 400
COLOR_PURPLE = RGBColor(168, 85, 247)    # Purple 500
COLOR_GREEN = RGBColor(34, 197, 94)      # Green 500
COLOR_TEXT_MAIN = RGBColor(255, 255, 255)# Crisp White
COLOR_TEXT_MUTED = RGBColor(148, 163, 184)# Slate 400
COLOR_TEXT_DIM = RGBColor(203, 213, 225) # Slate 300

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_slide_layout = prs.slide_layouts[6]


def create_base_slide(title_text: str, category_tag: str = "ENTERPRISE MCP GATEWAY"):
    slide = prs.slides.add_slide(blank_slide_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG
    bg.line.fill.background()

    # Category Pill Tag
    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(8.0), Inches(0.35))
    tf_tag = tag_box.text_frame
    tf_tag.word_wrap = True
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = category_tag.upper()
    p_tag.font.size = Pt(11)
    p_tag.font.bold = True
    p_tag.font.color.rgb = COLOR_ACCENT

    # Main Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.5), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_MAIN

    # Bottom Footer Bar
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.3))
    tf_footer = footer_box.text_frame
    p_foot = tf_footer.paragraphs[0]
    p_foot.text = "AI MCP Server Kit  |  Live Interactive Model Context Protocol Gateway  |  Confidential Customer Briefing"
    p_foot.font.size = Pt(10)
    p_foot.font.color.rgb = COLOR_TEXT_MUTED

    return slide


def add_card(slide, left: float, top: float, width: float, height: float, title: str, bullets: list, accent_color=COLOR_ACCENT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = COLOR_CARD_BORDER
    card.line.width = Pt(1.2)

    # Accent Top Stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + 0.05), Inches(top + 0.05), Inches(width - 0.1), Inches(0.08))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent_color
    stripe.line.fill.background()

    # Content Box
    tb = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.2), Inches(width - 0.5), Inches(height - 0.4))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(17)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_TEXT_MAIN
    p0.space_after = Pt(10)

    for b in bullets:
        p = tf.add_paragraph()
        p.text = f"•  {b}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_DIM
        p.space_after = Pt(6)


# ==========================================
# SLIDE 1: Title Slide (Hero)
# ==========================================
slide1 = prs.slides.add_slide(blank_slide_layout)
bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
bg1.fill.solid()
bg1.fill.fore_color.rgb = COLOR_BG
bg1.line.fill.background()

# Hero Badge
badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.3), Inches(4.2), Inches(0.45))
badge.fill.solid()
badge.fill.fore_color.rgb = RGBColor(30, 58, 138)
badge.line.color.rgb = COLOR_ACCENT
badge.line.width = Pt(1)
tf_b = badge.text_frame
p_b = tf_b.paragraphs[0]
p_b.text = "⚡ NEXT-GEN ENTERPRISE AI INFRASTRUCTURE"
p_b.font.size = Pt(11)
p_b.font.bold = True
p_b.font.color.rgb = COLOR_ACCENT
p_b.alignment = PP_ALIGN.CENTER

# Main Title & Subtitle
hero_tb = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(2.8))
tf_hero = hero_tb.text_frame
tf_hero.word_wrap = True

p_h1 = tf_hero.paragraphs[0]
p_h1.text = "AI MCP Server Kit & Gateway"
p_h1.font.size = Pt(40)
p_h1.font.bold = True
p_h1.font.color.rgb = COLOR_TEXT_MAIN

p_h2 = tf_hero.add_paragraph()
p_h2.text = "Secured, Dynamic Model Context Protocol Gateway & Live Architect"
p_h2.font.size = Pt(22)
p_h2.font.color.rgb = COLOR_ACCENT
p_h2.space_before = Pt(8)

p_h3 = tf_hero.add_paragraph()
p_h3.text = "Empowering AI Agents with Zero-Credential Leakage, Real-Time Schema Synthesis, and Enterprise Governance across 100+ Enterprise Systems."
p_h3.font.size = Pt(14)
p_h3.font.color.rgb = COLOR_TEXT_DIM
p_h3.space_before = Pt(16)

# Highlights Row
add_card(slide1, 1.2, 5.0, 3.4, 1.8, "🤖 Live LLM Architect", [
    "Mistral AI dynamic discovery",
    "Anti-hallucination validation",
    "15–45 rich tools per platform"
], COLOR_PURPLE)

add_card(slide1, 4.9, 5.0, 3.4, 1.8, "🔒 Zero Secret Exposure", [
    "Credentials isolated in .env",
    "Clean domain tool schemas",
    "Dual-tier Bearer auth"
], COLOR_GREEN)

add_card(slide1, 8.6, 5.0, 3.4, 1.8, "⚡ Multi-OS Gateway :5001", [
    "Multiplexed stdio execution",
    "Selective per-tool toggles",
    "1-Click PowerShell / CMD / Bash"
], COLOR_ACCENT)


# ==========================================
# SLIDE 2: The Enterprise Challenge
# ==========================================
slide2 = create_base_slide("The Enterprise Challenge: AI Agent Tool Sprawl & Security Risks", "PROBLEM STATEMENT")

add_card(slide2, 0.8, 1.8, 3.6, 4.8, "⚠️ Tool Sprawl & Fragmentation", [
    "Developers hand-code separate MCP servers for each tool (Jenkins, GitHub, ServiceNow, AWS, Terraform).",
    "Inconsistent schemas, mismatched parameters, and no standardized runtime.",
    "Days of manual coding per integration."
], RGBColor(239, 68, 68))

add_card(slide2, 4.8, 1.8, 3.6, 4.8, "🔓 Credential Leakage into LLMs", [
    "Naive MCP implementations pass API tokens and passwords directly inside LLM tool argument schemas.",
    "Critical enterprise secrets get logged into prompt transcripts and exposed to model contexts.",
    "Severe compliance and security risk."
], RGBColor(245, 158, 11))

add_card(slide2, 8.8, 1.8, 3.6, 4.8, "🛑 Governance & Invocation Chaos", [
    "No central gateway to enforce authentication, rate limits, or audit logging.",
    "Agents attempt to call unsupported or sensitive tools with no administrative whitelisting.",
    "OS command syntax confusion across Windows PowerShell, CMD, and Linux."
], RGBColor(236, 72, 153))


# ==========================================
# SLIDE 3: The Solution Overview
# ==========================================
slide3 = create_base_slide("The Solution: Unified MCP Gateway & Live Architect", "ENTERPRISE SOLUTION")

add_card(slide3, 0.8, 1.8, 3.6, 4.8, "1. Live Interactive Architect", [
    "Conversational chatbot on Port 5000 powered by Mistral AI LLM.",
    "Analyzes any platform request in real-time.",
    "Produces 15–45 comprehensive tool schemas with zero manual coding.",
    "Entity validation blocks fake/made-up tools."
], COLOR_PURPLE)

add_card(slide3, 4.8, 1.8, 3.6, 4.8, "2. Secured Gateway on Port 5001", [
    "Standardized JSON-RPC 2.0 endpoint: /mcp/<server_id>.",
    "Protected by Bearer Token & Gateway API Key.",
    "Multiplexes child servers over secure stdio pipes.",
    "Enforces selective tool exposure & toggling."
], COLOR_ACCENT)

add_card(slide3, 8.8, 1.8, 3.6, 4.8, "3. Complete Credential Vault", [
    "Secrets stored strictly server-side in mcp_servers/<server>/.env.",
    "LLM schemas only contain domain parameters (e.g. job_name, repo, incident_number).",
    "Zero risk of credential exfiltration."
], COLOR_GREEN)


# ==========================================
# SLIDE 4: Solution Architecture
# ==========================================
slide4 = create_base_slide("High-Level Enterprise Architecture & Data Flow", "ARCHITECTURE")

add_card(slide4, 0.8, 1.8, 3.6, 4.8, "Layer 1: AI Clients & UI", [
    "Web Chat UI (Port 5000)",
    "Autonomous AI Agents (LangChain, CrewAI, AutoGen, Claude Desktop)",
    "Developer Terminals (PowerShell, CMD, Bash)",
    "Standardized JSON-RPC 2.0 tools/call requests"
], COLOR_ACCENT)

add_card(slide4, 4.8, 1.8, 3.6, 4.8, "Layer 2: Core Gateway & Vault", [
    "Port 5001 Secured Gateway",
    "Bearer / API Key Authentication Filter",
    "Selective Tool Exposure Engine (Whitelist)",
    "Flexible Payload Parser (Windows/Linux)",
    "Keyring & Server-Side .env Vault"
], COLOR_PURPLE)

add_card(slide4, 8.8, 1.8, 3.6, 4.8, "Layer 3: Enterprise Connectors", [
    "Jenkins CI/CD (27 Real API Tools)",
    "GitHub Enterprise (45 Security & PR Tools)",
    "ServiceNow (Table API & Incident Ops)",
    "AWS Cloud (S3, EC2, IAM, Lambda)",
    "PostgreSQL, Jira, Terraform, Datadog"
], COLOR_GREEN)


# ==========================================
# SLIDE 5: Key Innovations & Differentiators
# ==========================================
slide5 = create_base_slide("Key Innovations & Architectural Differentiators", "INNOVATIONS")

add_card(slide5, 0.8, 1.8, 5.6, 2.3, "🛡️ Strict Server-Side Credential Isolation", [
    "Connection tokens/passwords never leak into tool function parameter schemas.",
    "Tools like get_jenkins_version() take empty {} args while authenticating internally."
], COLOR_GREEN)

add_card(slide5, 6.8, 1.8, 5.6, 2.3, "🧠 Anti-Hallucination Entity Validation", [
    "Mistral AI verifies if a platform exists before generating tools.",
    "Rejects nonsense words (e.g. 'bimbikili') and requests valid platform names."
], COLOR_PURPLE)

add_card(slide5, 0.8, 4.3, 5.6, 2.3, "☑️ Granular Tool Exposure Governance", [
    "Admins and developers can enable or disable individual tools with 1 click.",
    "Hidden tools are blocked at the Gateway level, preventing accidental execution."
], COLOR_ACCENT)

add_card(slide5, 6.8, 4.3, 5.6, 2.3, "💻 3-Tab Multi-OS Command Generator", [
    "Real-time executable cURL snippets tailored for PowerShell, CMD, and Linux Bash.",
    "One-click copy with full Bearer authentication and sample payload schemas."
], RGBColor(236, 72, 153))


# ==========================================
# SLIDE 6: Live Enterprise Connectors
# ==========================================
slide6 = create_base_slide("Live Enterprise Connectors in Action", "LIVE CONNECTORS")

add_card(slide6, 0.8, 1.8, 3.6, 4.8, "Jenkins CI/CD Suite (27 Tools)", [
    "• get_jenkins_version: Live controller status & version (2.541.2)",
    "• list_jobs: Real-time query of 33+ pipelines",
    "• get_job_details: Status & next build #",
    "• trigger_job_build: Parameterized execution",
    "• get_build_console_output: Log streaming",
    "• list_jenkins_nodes: Agent cluster fleet"
], COLOR_ACCENT)

add_card(slide6, 4.8, 1.8, 3.6, 4.8, "GitHub Enterprise (45 Tools)", [
    "• list_pull_requests: Filter by owner/repo",
    "• get_pr_status: CI checks & approvals",
    "• list_issues / create_issue: Issue tracker",
    "• list_code_scanning_alerts: SAST security",
    "• list_secret_scanning_alerts: Secret leaks",
    "• trigger_workflow: GitHub Actions trigger"
], COLOR_PURPLE)

add_card(slide6, 8.8, 1.8, 3.6, 4.8, "On-Demand Synthesizer", [
    "• ServiceNow: Incident, Change & CMDB ops",
    "• AWS S3: Bucket management & S3 objects",
    "• Terraform Cloud: Workspace & run execution",
    "• PostgreSQL DBA: Schema & query execution",
    "• Any REST/gRPC API in < 30 seconds"
], COLOR_GREEN)


# ==========================================
# SLIDE 7: Security & Governance
# ==========================================
slide7 = create_base_slide("Enterprise Security, Authentication & Governance", "SECURITY & COMPLIANCE")

add_card(slide7, 0.8, 1.8, 3.6, 4.8, "Gateway Perimeter Security", [
    "• Dedicated Port 5001 proxy with API Key verification.",
    "• Supports Authorization: Bearer <key> and X-API-Key headers.",
    "• Rejects unauthenticated requests with HTTP 401.",
    "• Fully customizable key management in UI."
], COLOR_ACCENT)

add_card(slide7, 4.8, 1.8, 3.6, 4.8, "Credential Vault Isolation", [
    "• Per-server credentials stored in isolated .env files.",
    "• OS Keyring integration (Windows Credential Manager / Linux SecretService).",
    "• Secrets are injected at runtime into stdio subprocesses.",
    "• Zero credential transmission to LLM clients."
], COLOR_GREEN)

add_card(slide7, 8.8, 1.8, 3.6, 4.8, "Auditability & Policy Controls", [
    "• Detailed Gateway access logging with timestamps.",
    "• Selective tool exposure enables least-privilege access for autonomous agents.",
    "• Ready for enterprise SIEM / Langfuse observability integration."
], COLOR_PURPLE)


# ==========================================
# SLIDE 8: Developer & Agent Experience
# ==========================================
slide8 = create_base_slide("Seamless Developer & Agent Experience", "INTEGRATION")

add_card(slide8, 0.8, 1.8, 5.6, 4.8, "Interactive Web Dashboard (Port 5000)", [
    "• Left Rail: Connected MCP servers with real-time status badges.",
    "• Server Details Canvas: Gateway endpoint URL, credentials location, transport details.",
    "• Function Catalog: Displays parameter types (string, integer, object) and REQUIRED / OPTIONAL badges.",
    "• Interactive Checkboxes: Instant toggle to expose or hide specific functions.",
    "• Chat Architect Modal: Live multi-turn conversation with Mistral AI."
], COLOR_ACCENT)

add_card(slide8, 6.8, 1.8, 5.6, 4.8, "Universal Agent & CLI Compatibility", [
    "• PowerShell (Windows): Backtick continuation, JSON escape handling.",
    "• CMD (Windows): Caret continuation with escaped quotes.",
    "• Linux / macOS (Bash): Standard POSIX cURL with formatted payloads.",
    "• AI Framework Ready: 1-line integration with LangChain, LlamaIndex, CrewAI, AutoGen, and Claude Desktop config.json."
], COLOR_PURPLE)


# ==========================================
# SLIDE 9: Business ROI & Value
# ==========================================
slide9 = create_base_slide("Business Impact, ROI & Time-to-Value", "BUSINESS VALUE")

add_card(slide9, 0.8, 1.8, 3.6, 4.8, "⏱️ 95% Faster Time-to-Market", [
    "• Traditional: 3–5 days per custom MCP server integration.",
    "• With AI MCP Kit: Synthesize and deploy comprehensive 25-tool suites in < 60 seconds.",
    "• Accelerates enterprise AI agent rollout by months."
], COLOR_GREEN)

add_card(slide9, 4.8, 1.8, 3.6, 4.8, "🛡️ 100% Secret Protection", [
    "• Eliminates accidental credential exposure in LLM prompts.",
    "• Compliant with SOC2, ISO 27001, and enterprise security policies.",
    "• Centralized key rotation and access revocation."
], COLOR_ACCENT)

add_card(slide9, 8.8, 1.8, 3.6, 4.8, "📈 Massive Scalability", [
    "• Single unified gateway multiplexes dozens of micro-servers.",
    "• Zero vendor lock-in; open Model Context Protocol standard.",
    "• Connects legacy internal tools as easily as modern cloud APIs."
], COLOR_PURPLE)


# ==========================================
# SLIDE 10: Summary & Next Steps
# ==========================================
slide10 = create_base_slide("Conclusion & Live Demonstration", "NEXT STEPS")

add_card(slide10, 0.8, 1.8, 5.6, 4.8, "Summary of Delivered Value", [
    "✅ Unified MCP Gateway running on Port 5001 with Bearer authentication.",
    "✅ Live Mistral AI Architect with Anti-Hallucination validation.",
    "✅ Strict server-side credential isolation in .env vault.",
    "✅ Live Jenkins (27 tools) and GitHub Enterprise (45 tools) connectors verified.",
    "✅ Interactive UI with per-tool exposure checkboxes and 3-tab multi-OS command generator."
], COLOR_GREEN)

add_card(slide10, 6.8, 1.8, 5.6, 4.8, "Demonstration Agenda", [
    "1. Live Chat Demo: Requesting new platform with Mistral AI (Port 5000).",
    "2. Anti-Hallucination Check: Testing invalid inputs.",
    "3. Checkbox Governance: Toggling tool exposure on the fly.",
    "4. Live Execution: Running Jenkins & GitHub tools via PowerShell cURL against Port 5001.",
    "5. Q&A and Deployment Discussion."
], COLOR_ACCENT)

prs.save(OUTPUT_PPTX)
print(f"Successfully generated presentation deck at: {OUTPUT_PPTX}")
